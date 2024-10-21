import re
import pymupdf
from nltk.tokenize import sent_tokenize
from docx.api import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pypandoc

def clean_text(text):
    text = re.sub(r"[^\x00-\x7F]", " ", text)
    text = re.sub(r"[\n]", " ", text)
    text = re.sub(r'\s+', ' ', text).strip()
    text = re.sub(r'[“”]', '"', text)
    text = re.sub(r"[‘’]", "'", text)
    text = text.replace('\xad', '')
    text = re.sub(r'[‒–—―]', '-', text)
    return text

# Function to create text chunks
def segment_text(text, max_segment_length=700, batch_size=7):
    sentences = sent_tokenize(text)
    segments = []
    current_segment = ""
    
    for sentence in sentences:
        if len(current_segment) + len(sentence) <= max_segment_length:
            current_segment += sentence + " "
        else:
            segments.append(current_segment.strip())
            current_segment = sentence + " "
    
    if current_segment:
        segments.append(current_segment.strip())
    
    # Create batches
    batches = [segments[i:i + batch_size] for i in range(0, len(segments), batch_size)]
    return batches 

def get_pdf_text(pdf_file):
    doc = pymupdf.open(stream=pdf_file.read(), filetype="pdf")
    text = ""
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text

# Function to get text from a DOCX file
def get_doc_text(doc_files):
    doc = Document(doc_files)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to get text from a PPTX file
def get_ppt_text(ppt_files):
    prs = Presentation(ppt_files)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to get text from HTML files
def get_html_text(html_files):
    text = ""
    for html_file in html_files:
        with open(html_file, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            text += soup.get_text()
    return text

# Function to get text from LaTeX files
def get_latex_text(latex_files):
    text = ""
    for latex_file in latex_files:
        output = pypandoc.convert_file(latex_file, 'plain')
        text += output
    return text

# Function to parse text from a file
def parse_text(file):
    text = file.getvalue().decode("utf-8")
    return text

# Function to get text from uploaded documents
def get_text_from_document(file):
    content = ""
    if file is not None:
        if file.name.endswith('.pdf'):
            content += get_pdf_text(file)
        elif file.name.endswith('.docx') or file.name.endswith('.doc'):
            content += get_doc_text(file)
        elif file.name.endswith('.pptx') or file.name.endswith('.ppt'):
            content += get_ppt_text(file)
        elif file.name.endswith('.html'):
            content += get_html_text(file)
        elif file.name.endswith('.tex'):
            content += get_latex_text(file)
        elif file.name.endswith('.txt'):
            content += parse_text(file)
    return content
